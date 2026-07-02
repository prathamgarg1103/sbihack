# Finalizes Diya-Deck-v2.pptx: swaps the 8 labeled SCREENSHOT/QR placeholders for real
# captures (cropped to fit each box), fills in the team name, and updates the stale
# test count. Re-runnable only against a fresh deck (placeholders must still exist).
from pathlib import Path

import qrcode
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches

ROOT = Path(r"E:\sbi hack")
CAP = ROOT / "video" / "captures"
DECK_IMG = CAP / "deck"
DECK_IMG.mkdir(exist_ok=True)

LIVE_URL = "https://saarthi-nu.vercel.app"
TEAM = "Team Diya — Pratham Garg · Thapar Institute of Engineering & Technology"


def crop(src: str, out: str, top: float = 0.0, bottom: float = 1.0) -> Path:
    im = Image.open(CAP / src)
    w, h = im.size
    im.crop((0, int(h * top), w, int(h * bottom))).save(DECK_IMG / out)
    return DECK_IMG / out


# --- prepare images, cropped to suit each placeholder box -----------------------
images = {
    "s6": crop("e02d_comparison.png", "deck_s6.png", 0.0, 0.50),   # 3 columns + first rows
    "s7": crop("e03c_exit_vs_stay.png", "deck_s7.png", 0.0, 0.72),  # exit-now + stay panels
    "s8": crop("e02c_why_da.png", "deck_s8.png", 0.42, 1.0),        # why card + DA objection
    "s9": crop("e07_mission_control.png", "deck_s9.png", 0.105, 0.39),  # status + feed rows
    "s10a": crop("e02b_nudge_hi.png", "deck_s10a.png", 0.55, 1.0),  # Hindi nudge card
    "s10b": crop("e06b_guardian_phone.png", "deck_s10b.png"),       # guardian approval card
    "s12": crop("e08b_compliance_top.png", "deck_s12.png", 0.0, 0.50),  # chain-verified badge
}
qr = qrcode.make(LIVE_URL)
qr_path = DECK_IMG / "deck_qr.png"
qr.save(qr_path)
images["qr"] = qr_path

# (slide_number, match_text, image_key)
TARGETS = [
    (6, "SCREENSHOT", "s6"),
    (7, "SCREENSHOT", "s7"),
    (8, "SCREENSHOT", "s8"),
    (9, "SCREENSHOT", "s9"),
    (10, "Hindi nudge", "s10a"),
    (10, "Guardian co-consent", "s10b"),
    (12, "SCREENSHOT", "s12"),
    (14, "QR CODE", "qr"),
]

prs = Presentation(ROOT / "Diya-Deck-v2.pptx")
done = []

for slide_no, match, key in TARGETS:
    slide = prs.slides[slide_no - 1]
    label = None
    for sh in slide.shapes:
        if sh.has_text_frame and match in sh.text_frame.text and (
            "SCREENSHOT" in sh.text_frame.text or "QR CODE" in sh.text_frame.text
        ):
            label = sh
            break
    if label is None:
        print(f"!! slide {slide_no}: no placeholder matching {match!r}")
        continue
    box_l, box_t = Emu(label.left).inches, Emu(label.top).inches
    box_w, box_h = Emu(label.width).inches, Emu(label.height).inches
    im = Image.open(images[key])
    iw, ih = im.size
    scale = min(box_w / iw, box_h / ih) * 0.97  # slight inset inside the frame
    w, h = iw * scale, ih * scale
    left = box_l + (box_w - w) / 2
    top = box_t + (box_h - h) / 2
    slide.shapes.add_picture(str(images[key]), Inches(left), Inches(top), Inches(w), Inches(h))
    label._element.getparent().remove(label._element)
    done.append(f"slide {slide_no}: {key} ({w:.2f}x{h:.2f} in)")

# --- team name on slide 1 --------------------------------------------------------
for sh in prs.slides[0].shapes:
    if sh.has_text_frame and "[Team name" in sh.text_frame.text:
        for para in sh.text_frame.paragraphs:
            if "[Team name" in "".join(r.text for r in para.runs):
                para.runs[0].text = TEAM
                for r in list(para.runs[1:]):
                    r._r.getparent().remove(r._r)
                done.append("slide 1: team name filled")
                break

# --- stale test count on slide 11 -------------------------------------------------
for sh in prs.slides[10].shapes:
    if sh.has_text_frame and "26 tests" in sh.text_frame.text:
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if "26 tests" in r.text:
                    r.text = r.text.replace("26 tests", "44 tests")
                    done.append("slide 11: 26 -> 44 tests")

prs.save(ROOT / "Diya-Deck-v2.pptx")
print("\n".join(done))
print(f"\nremaining placeholders: ", end="")
prs2 = Presentation(ROOT / "Diya-Deck-v2.pptx")
left = [
    f"slide {i}"
    for i, s in enumerate(prs2.slides, 1)
    for sh in s.shapes
    if sh.has_text_frame and ("SCREENSHOT" in sh.text_frame.text or "QR CODE" in sh.text_frame.text)
    and i != 15  # slide 15 is the hidden inventory appendix; its mentions are intentional
]
print(", ".join(left) or "none (outside hidden appendix)")
